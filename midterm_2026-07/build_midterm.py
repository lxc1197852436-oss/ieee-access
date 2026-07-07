# -*- coding: utf-8 -*-
import os
import sys
from docx import Document
from docx.shared import Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT_DIR = r'd:\ieee access\ieee-access\midterm_2026-07'
os.makedirs(OUT_DIR, exist_ok=True)

DARK_BLUE = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x0E, 0x7C, 0x86)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)
CN = '宋体'
CNH = '黑体'

# 所有正文内容用单引号包裹；中文引号统一用中文全角〔〕替代以彻底避免冲突，
# 生成后再无影响。这里直接用纯文本，内部不出现 ASCII 单引号。

def set_cn(run, font=CN):
    run.font.name = font
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn('w:rFonts'))
    if rFonts is None:
        rFonts = rPr.makeelement(qn('w:rFonts'), {})
        rPr.append(rFonts)
    rFonts.set(qn('w:eastAsia'), font)

def H(doc, text, level):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(12)
    p.paragraph_format.space_after = Pt(6)
    r = p.add_run(text)
    if level == 0:
        r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = DARK_BLUE
        set_cn(r, CNH); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    elif level == 1:
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = DARK_BLUE; set_cn(r, CNH)
    else:
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT; set_cn(r, CNH)
    return p

def B(doc, text, indent=True):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.line_spacing = 1.5
    if indent:
        p.paragraph_format.first_line_indent = Cm(0.74)
    r = p.add_run(text)
    r.font.size = Pt(12); set_cn(r, CN)
    return p

def BL(doc, text):
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Cm(0.74)
    p.paragraph_format.space_after = Pt(2)
    p.paragraph_format.line_spacing = 1.4
    r = p.add_run('• ' + text)
    r.font.size = Pt(12); set_cn(r, CN)
    return p

def EQ(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(text)
    r.font.size = Pt(12); r.font.italic = True
    return p

def TBL(doc, header, rows):
    t = doc.add_table(rows=1+len(rows), cols=len(header))
    t.style = 'Light Grid Accent 1'
    for i, h in enumerate(header):
        c = t.rows[0].cells[i]; c.text = ''
        r = c.paragraphs[0].add_run(h)
        r.font.bold = True; r.font.size = Pt(11); r.font.color.rgb = WHITE; set_cn(r, CNH)
        tcPr = c._tc.get_or_add_tcPr()
        shd = tcPr.makeelement(qn('w:shd'), {qn('w:fill'): '1F3A5F'})
        tcPr.append(shd)
    for ri, row in enumerate(rows):
        for i, v in enumerate(row):
            c = t.rows[ri+1].cells[i]; c.text = ''
            r = c.paragraphs[0].add_run(str(v))
            r.font.size = Pt(11); set_cn(r, CN)
    return t

def build_docx():
    doc = Document()
    for s in doc.sections:
        s.top_margin = Cm(2.5); s.bottom_margin = Cm(2.5)
        s.left_margin = Cm(3.0); s.right_margin = Cm(3.0)
    H(doc, '研究生中期检查报告', 0)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for label, val in [('课题名称：', '大语言模型与深度强化学习融合的虚拟电厂动态优化与决策研究'),
                       ('研究生：', '李星辰'), ('导　　师：', '王继军'),
                       ('学　　校：', '上海电力大学'), ('报告日期：', '2026 年 7 月')]:
        rp = p.add_run(label + val + '\n'); rp.font.size = Pt(13); set_cn(rp, CN)
    doc.add_paragraph()

    H(doc, '一、课题研究进展概述', 1)
    B(doc, '本课题围绕大语言模型与深度强化学习融合的虚拟电厂动态优化与决策展开研究。截至中期检查节点，研究已从开题阶段的理论设计推进到方法实现、长训练实验、多场景评估与新事件自适应验证的闭环阶段，并已形成一篇投稿至 IEEE Access 的英文论文初稿（题名：Event-Aware Virtual Power Plant Dispatch via Language-Enhanced Soft Actor-Critic Reinforcement Learning With a Semantic Safety Layer）。')
    B(doc, '与开题时设定的语义启发式策略加离散 Soft-Q 基线过渡方案相比，中期阶段已完成两步关键升级：第一，将决策核心从离散动作 Soft-Q 升级为连续动作 Soft Actor-Critic（SAC），可直接输出储能充放电功率；第二，将大语言模型的角色明确界定为语义编码器（DeepSeek）而非最终决策器，并提出语义安全层加事件覆盖相关权重的混合控制结构。在此基础上，进一步完成了新事件（负电价深过剩）自适应研究、参数变体稳健性验证、德国 DE-LU 真实电价与 Open-Meteo 真实气象验证，以及基于 Q 差异的学习型事件覆盖门控。上述工作构成了中期答辩的阶段性成果。')
    B(doc, '需要特别说明的是，本研究所用 15 分钟级虚拟电厂调度数据为国内公开数据校准的仿真序列，而非真实虚拟电厂运行明细。这一数据边界在论文与报告中均作明确界定，以避免夸大公开数据的真实性。')

    H(doc, '二、研究背景与意义', 1)
    B(doc, '在双碳目标与新型电力系统建设背景下，分布式光伏、用户侧储能、可调负荷、电动汽车等灵活性资源快速接入配电侧与用户侧。虚拟电厂（VPP）通过信息通信、聚合控制与优化调度，将这些分散资源组织为可观、可测、可控、可交易的整体，为新能源消纳、需求响应、辅助服务和电力市场互动提供支撑。然而，VPP 运行具有显著的动态性、不确定性与多源信息融合特征：一方面，负荷、光伏出力、电价与储能荷电状态（SOC）持续变化；另一方面，气象预警、市场公告、需求响应邀约、新能源消纳提示等外部信息通常以自然语言文本形式出现，且往往先于其数值效应反映在时间序列中。')
    B(doc, '现有强化学习调度方法多仅消费数值状态，对自然语言事件利用不足，导致策略在高温负荷压力、价格尖峰、新能源消纳压力等场景下缺少前瞻性。大语言模型（LLM）为非结构化文本与数值决策模型之间提供了接口，但直接让 LLM 输出最终调度动作存在可审计性、可复现性与物理可行性风险。因此，本课题采用解耦式语言增强设计：LLM 仅将事件文本转化为结构化风险特征，最终动作仍由可复现的 SAC 策略与确定性安全层输出。这一设计兼顾了文本理解能力、调度可复现性与实验可评价性，对推动 LLM 在电力调度中的落地具有方法学意义。')
    B(doc, '从国内政策与市场环境看，2024 年国家发改委、国家能源局相继出台虚拟电厂相关指导意见，明确将虚拟电厂作为新型电力系统负荷管理与社会资本参与电力市场的重要载体；广东、山东、山西等现货市场连续运行，独立储能、聚合商等主体参与电能量与辅助服务市场的机制逐步完善。这使得 VPP 调度不仅要应对物理层面的源荷波动，还要在市场层面响应价格信号与事件邀约。文本事件（如需求响应邀约、负电价预警、消纳受限通知）在市场出清前后均会发布，蕴含了数值时间序列尚未反映的前瞻信息。能否让调度策略理解并利用这类信息，直接关系到 VPP 在事件场景下的经济性与风险表现，这正是本课题的切入动机。')

    H(doc, '三、已完成研究工作', 1)
    H(doc, '3.1 公开数据校准的虚拟电厂数据管线', 2)
    B(doc, '针对国内真实虚拟电厂 15 分钟级运行明细与完整现货出清数据通常不公开的现实，本研究采用公开数据校准加仿真构造的方式建立可复现实验环境。已完成的数据工作包括：')
    BL(doc, '国家能源局 2024 年全社会用电量（98521 亿千瓦时，同比增长 6.8%）与分产业用电量数据；')
    BL(doc, '国家能源局 2024 年可再生能源装机与发电量数据（可再生能源累计装机 18.89 亿千瓦，风电 5.21 亿千瓦、太阳能 8.87 亿千瓦）；')
    BL(doc, '广东电力现货市场与独立储能公开披露数据（价格区间、充放电电量、充放电均价与价差）；')
    BL(doc, '基于上述公开数据校准的广东虚拟电厂 15 分钟样例场景，覆盖负荷、光伏、电价、温度、SOC 与中文文本事件。')
    B(doc, '论文中对数据边界作了显式界定：公开披露数据用于校准研究背景、价格范围与运行参数；强化学习状态轨迹为可复现的仿真生成序列，而非真实私有遥测。这一界定使实验可复现，同时避免对数据真实性作过度声明。')

    H(doc, '3.2 虚拟电厂强化学习环境', 2)
    B(doc, '将 VPP 调度建模为马尔可夫决策过程。数值状态包含负荷、光伏出力、电价、温度、SOC 与时间编码（正余弦小时特征）；动作为储能连续充放电功率（正为放电、负为充电）；环境按 SOC 与功率约束裁剪不可行动作。奖励函数统一为：')
    EQ(doc, 'r_t = q_t(net)·Δt·π_t − c_deg·|a_t|·Δt − c_cur·q_t(cur) − c_vio·I_t')
    B(doc, '其中 q_t(net) 为动作后的净交换电量，c_deg 为退化成本，q_t(cur) 为弃光量，I_t 为不可行指示。训练采用相对无动作基线的优势型奖励，测试报告原始运行收益。')

    H(doc, '3.3 语言增强 SAC 框架（LE-DRL-SAC）', 2)
    B(doc, '提出 LE-DRL-SAC 框架，其核心为解耦式语言增强：DeepSeek 大语言模型作为语义编码器，对每个事件模板查询一次并缓存结果，将中文事件文本转化为五维结构化语义向量（总体风险、价格尖峰压力、负荷压力、新能源消纳压力、储能偏置），与数值状态拼接形成增强状态 s_aug=[s_num, s_sem]。LLM 不参与逐步控制回路，不引入在线推理延迟，亦不输出调度动作。')
    B(doc, '在 SAC 最大熵目标之上，引入 Actor 语义一致性正则：')
    EQ(doc, 'J_π(LE) = J_π + λ_sem · E[‖a_t − a_sem(s_aug)‖²]')
    B(doc, '其中 a_sem(·) 为由价格、SOC 与风险特征导出的参考动作方向。该正则在训练阶段促使 Actor 学习语义特征对应的合理储能行为。')

    H(doc, '3.4 语义安全层与事件覆盖相关权重', 2)
    B(doc, '在测试阶段，将同一风险一致先验用作确定性语义安全层，与学习动作进行有界混合：')
    EQ(doc, 'a_final = (1−w)·a_SAC + w·a_sem')
    B(doc, '最终动作再经储能功率与 SOC 可行性裁剪。关键设计在于：权重 w 被视为事件覆盖相关量，而非通用常数。在先验所设计的事件类别上，先验接近场景最优，取高 w；在先验未覆盖的事件类别上，先验无对应分支，取低 w 让学习 Actor 接管。在四个已知事件场景的主实验中，经权重扫描选定 w=0.9。为隔离文本语义贡献与安全层机制本身，另设数值安全层消融：用仅基于价格、SOC、小时与光伏过剩的纯数值先验替代语义先验，采用相同混合公式。')
    B(doc, '需要强调一项设计权衡：为何不让 LLM 直接充当决策器？一方面，VPP 储能调度必须满足 SOC 动态、功率边界与市场结算逻辑，LLM 自由生成的动作难以保证逐可行；另一方面，调度结果需要可复现、可审计，以便事后核查与责任界定。本框架将 LLM 限定在每事件模板查询一次并缓存的离线语义编码角色，既利用其文本理解能力，又使其不进入逐步控制回路，从而把不可控性隔离在缓存层之外。这一处理方式与电力领域 LLM 综述中所倡导的 LLM 作特征抽取器或顾问而非最终控制器思路一致，但本文进一步将其落到连续动作储能调度的完整闭环中，并以数值安全层消融提供了机制层面的对照。')

    H(doc, '3.5 新事件自适应与真实电价验证', 2)
    B(doc, '为检验框架对先验未覆盖事件类别的适应能力，构造第五类场景 S5（负电价深过剩）：午间光伏过剩压低现货价格至负，晚高峰价格受抑而非尖峰。先验的晚高峰放电源分支在 S5 上失配，而 DeepSeek 对该新事件文本返回消纳压力 0.90、储能偏置 +0.80，正确识别其过剩本质。')
    B(doc, '在 S5 上采用松弛正则化变体（正则权重置零、安全层 w=0），Actor 在 S1–S5 上重训练（5 个随机种子、80 episode），与同条件重训练的 SAC-Numeric 对比。进一步构造四个结构不同的负电价变体 V1–V4（改变负电价深度、季节、光伏装机），验证适应性是否稳健。最重要的是，构造了完全由真实德国 DE-LU 日前价格（2026-07-04，15 分钟分辨率，27 个负电价区间，最低 −11.8 EUR/MWh）与真实 Open-Meteo 气象组成的验证场景，以检验方法在真实浅负电价下是否仍有效。')

    H(doc, '3.6 学习型事件覆盖门控', 2)
    B(doc, '两套固定配置（已知事件高 w、新事件 w=0 松弛正则）需在部署时已知覆盖标签。为检验该标签能否从数据中学习，构造门控混合专家控制器：两个 SAC 专家共享同一 LLM 增强状态，分别带或不带语义一致性正则；一个小门控网络 g_ψ(s)→w∈[0,1] 以逐批归一化的 Q 差异为二分类目标训练，无需覆盖标签，输出即为混合权重。')

    H(doc, '3.7 可复现实验系统与展示平台', 2)
    B(doc, '基于 Python 实现数据处理、VPP 环境、SAC 训练、指标计算与结果导出；基于 FastAPI 与本地 Dashboard 构建展示系统，支持场景生成、策略运行、收益与 CVaR 指标展示、SOC 与动作轨迹可视化，以及中文文本事件与 AI 解释结果展示。')

    H(doc, '四、阶段性实验结果', 1)
    B(doc, '实验设置四类七日场景：S1 常规夏季运行、S2 高温负荷压力、S3 价格尖峰、S4 新能源消纳压力。SAC-Numeric、LE-DRL w/o Text、LE-DRL-SAC 以 80 episode、种子 2026/2031/2042 训练；新事件与参数变体研究扩展至 5 个种子。')

    H(doc, '4.1 跨场景基线对比', 2)
    B(doc, '在 S1–S4 跨场景平均总收益上，最终控制器 LE-DRL-SAC 加语义安全层（w=0.9）取得 −208,468.9 元，优于 SAC-Numeric（−211,023.2 元）、Rule-Based（−208,969.9 元）、增强滚动时域（−209,988.3 元）、Linear-MPC（−210,656.0 元）及数值安全层 SAC（−209,125.3 元）。采用种子配对 bootstrap（12 个匹配的种子场景对，差值重采样 20,000 次）：相对 SAC-Numeric 配对差 +2,554.2 元，95% CI [+2,060.4, +3,043.8]，Wilcoxon p=0.0005。')
    B(doc, '文本消融显示 SAC-Numeric 与 LE-DRL w/o Text 几乎完全一致（−211,023.2 对 −211,023.1 元），而 LE-DRL-SAC 在注入非零语义分数后改善，说明增益并非仅来自输入维度增大。数值安全层将 SAC-Numeric 从 −211,023.2 提升至 −209,125.3 元，证明安全层机制本身有价值；语义安全层进一步达到 −208,468.9 元，说明剩余改善可归因于结构化文本语义而非混合机制本身。')
    TBL(doc, ['模型', '总收益(元)↑', 'CVaR 5%(元)', '吞吐(MWh)', '高价放电率'],
        [['LE-DRL-SAC+安全层(w=0.9)', '−208,468.9', '−696.7', '28.37', '0.231'],
         ['Rule-Based', '−208,969.9', '−705.0', '19.94', '0.176'],
         ['SAC-Numeric+数值安全层', '−209,125.3', '−704.0', '16.92', '0.193'],
         ['增强滚动时域', '−209,988.3', '−713.8', '14.80', '0.095'],
         ['Linear-MPC', '−210,656.0', '−714.6', '10.81', '0.073'],
         ['SAC-Numeric', '−211,023.2', '−721.3', '3.04', '0.000'],
         ['LE-DRL w/o Text', '−211,023.1', '−721.3', '3.04', '0.000']])

    H(doc, '4.2 新事件自适应（S5）', 2)
    B(doc, '松弛正则化 LE-DRL-SAC 在 S5 取得 −141,068.4 元，优于同条件重训练的 SAC-Numeric（−146,550.2 元）5,481.7 元，5 种子配对 bootstrap 95% CI [+4,875.3, +6,074.0]，5/5 种子有利。行为指标证实增益来自事件自适应：松弛版在 100% 负电价步充电，并学会在晚高峰价格受抑日避免先验的晚高峰放电。')
    B(doc, '正则敏感性扫描揭示 Actor 语义一致性正则的双重角色：在四个已知事件上收益跨权重近乎持平（0.5% 以内），但在 S5 上随正则松弛单调改善，从满权重 −159,670.6 元改善至零权重 −141,155.6 元，增益 18,515 元，且已知事件不退化。这表明语义信息有两种可分离作用——作为 Actor 可学习的状态特征，与作为约束 Actor 的训练目标——二者在新事件上冲突。')

    H(doc, '4.3 参数变体稳健性与真实电价验证', 2)
    B(doc, '四个负电价变体 V1/V3/V4 均给出显著正间隙（95% CI 全在零以上，5/5 种子有利），V2 趋正。跨深度、季节、光伏装机三维度的稳健性表明，新事件自适应非单一参数设定的产物。')
    B(doc, '真实电价验证（最接近真实市场的证据）：在德国 DE-LU 真实日前价格加真实 Open-Meteo 气象上，松弛正则化 LE-DRL-SAC 取得 −34,955.7 EUR，优于 SAC-Numeric 的 −37,156.3 EUR，间隙 +2,200.6 EUR，5 种子配对 bootstrap 95% CI [+593.4, +3,627.3]，4/5 种子有利。这表明新事件自适应可迁移至真实负电价记录与真实气象，而非仅限仿真生成器。')
    TBL(doc, ['场景', '间隙(EUR)', '95% CI (EUR)', '有利种子'],
        [['真实 DE-LU 价格+真实气象', '+2,201', '[+593, +3,627]', '4/5']])

    H(doc, '4.4 学习型事件覆盖门控', 2)
    B(doc, '门控在 S1–S5（5 种子）上训练。在四个已知事件上平均门控权重 w=0.65（偏向正则化先验专家），在 S5 与四个同结构变体 V1–V4 上 w 降至 0.31–0.48（偏向自由专家）。这表明仅以 Q 差异为信号、不告知事件覆盖标签，学习型门控即可恢复已知事件高 w、负电价事件低 w 的分割，验证了事件覆盖相关权重可从数据中自动化学习。')

    H(doc, '4.5 结果综合讨论', 2)
    B(doc, '综合上述实验可得出四点判断。其一，结构化文本语义特征对 VPP 储能调度确有增益，且该增益可被文本消融与数值安全层消融双重隔离，排除了输入维度增大与混合机制本身两类混淆解释。其二，增益的大小与场景高度相关：在事件驱动场景（S2 高温负荷压力、S4 消纳压力、S5 负电价）上更明显，在常规运行场景（S1）上相对收窄，这与文本语义在事件场景下价值最大的直觉一致，也提示后续应把评估重心放在事件场景。其三，Actor 语义一致性正则的两重作用是本研究最重要的方法学发现——它既是已知事件上的训练稳定器，也是新事件上的适应约束，二者的冲突决定了权重不能取通用常数，而应随事件覆盖调整。其四，真实电价与真实气象验证是全部证据中最接近部署的一环，它在远比仿真 S5 浅的真实负电价下仍观察到显著正间隙，说明事件感知机制并非仿真生成器的产物。需要客观指出的是，跨场景平均收益的绝对百分比改善幅度有限（相对 Rule-Based 约 0.24%），这是奖励结构本身决定的——总收益的大部分来自任何可行策略都能获取的净出口收入，储能控制的边际价值相对总收益天然较小；但 CVaR 5% 与高价放电率等行为指标显示，策略在尾部风险抑制与高价时段放电集中度上有更具操作意义的改善。')

    H(doc, '五、当前主要创新点', 1)
    BL(doc, '解耦式语言增强 SAC 框架：将 LLM 严格界定为语义编码器（每事件模板查询一次并缓存），与连续动作 SAC 结合，LLM 不输出调度动作，最终决策确定、可审计、物理可行。')
    BL(doc, '事件覆盖相关的语义安全层：提出权重 w 视事件覆盖而定的混合控制原则，区别于通用常数权重；通过数值安全层消融隔离文本语义贡献与机制贡献。')
    BL(doc, '新事件自适应研究：引入先验未覆盖的负电价事件类别，证明松弛正则化下 LLM 特征可使 SAC 适应先验无法覆盖的事件（+5,482 元，5/5 种子），并区分语义信息的两种可分离作用。')
    BL(doc, '真实电价与真实气象验证：在德国 DE-LU 真实日前价格与 Open-Meteo 真实气象上验证自适应迁移（+2,201 EUR，4/5 种子），为标题事件感知提供最贴近部署的实证锚点。')
    BL(doc, '学习型事件覆盖门控：以逐批归一化 Q 差异为无标签信号，从数据中恢复事件覆盖分割，验证覆盖相关权重可自动化。')

    H(doc, '六、存在问题与不足', 1)
    BL(doc, '数据真实性边界：主跨场景轨迹为公开数据校准的仿真序列，非真实私有 VPP 遥测；OOD 真实气象测试与德国真实电价验证已部分缓解，但仍需在实测 VPP/微电网遥测上验证。')
    BL(doc, 'DeepSeek 相对简单编码器的独特性待加强：当前证据支持结构化语义特征优于零文本与数值先验消融，但 DeepSeek 相对关键词编码器的独特优势尚需专用对比坐实。若该对比显示 DeepSeek 与关键词编码器统计不可区分，则需将语言增强的论据重心从 LLM 独特理解力调整为 LLM 作为获取结构化语义特征的鲁棒自动化手段（无需手工关键词），并强化 native 优于 w/o Text、noisy 劣于 w/o Text 两条坐实证据。')
    BL(doc, '统计方法已修正：早期草稿基于 3 个种子均值的 bootstrap 给出过窄 CI（宽度仅 55 元），现已改为种子配对 bootstrap 加 Wilcoxon 非参交叉检验，反映场景级变异性。')
    BL(doc, '优化基线强度：当前含 Linear-MPC 与滚动时域，但完整 MILP/MPC 滚动优化器在相同预测与储能约束下的对比仍需补充。')
    BL(doc, '门控优化：当前门控仅演示覆盖分割可学习，门控混合动作的奖励最优化有待后续。')

    H(doc, '七、下一步研究计划', 1)
    B(doc, '中期答辩定于 2026 年 7 月 9 日。下一阶段工作按以下时间节点推进：')
    TBL(doc, ['阶段', '时间', '主要工作'],
        [['中期答辩', '2026 年 7 月', '完成中期检查报告与答辩，汇总专家意见'],
         ['实验完善', '2026 年 8–10 月', '补充 DeepSeek vs 关键词/noisy 专用消融；补完整 MILP/MPC 基线；扩充场景与种子；门控奖励最优化'],
         ['论文定稿', '2026 年 11–12 月', '完成 IEEE Access 投稿论文修订与补充实验；同步推进学位论文各章定稿'],
         ['预答辩与送审', '2027 年 1–3 月', '学位论文全稿查重、预答辩、外审修改'],
         ['答辩', '2027 年上半年', '正式学位论文答辩']])
    B(doc, '具体研究层面下一步重点：')
    BL(doc, '强化 LLM 独特性论证：补 DeepSeek 与关键词编码器、noisy 分数的专用对比，并探索恶意关键词消融以坐实 LLM 理解力优势。')
    BL(doc, '补强优化基线：在相同预测与储能约束下实现完整 MILP 与 MPC 滚动优化器对比。')
    BL(doc, '实测数据验证：争取在实测 VPP/微电网遥测或公开实测数据集上验证方法。')
    BL(doc, '门控与安全约束：将门控扩展为奖励最优的混合动作选择，并引入形式化安全约束 RL 以增强可部署性。')
    BL(doc, '学位论文写作：按七章结构推进各章定稿，统一图表与公式编号、参考文献格式。')

    H(doc, '八、中期答辩建议表述', 1)
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Cm(0.74)
    p.paragraph_format.line_spacing = 1.5
    r = p.add_run('本研究已完成面向中国虚拟电厂场景的公开数据校准数据管线、VPP 强化学习环境、语言增强 SAC 框架（LE-DRL-SAC）、语义安全层、新事件自适应研究、真实电价验证与学习型事件覆盖门控。实验表明，结构化文本语义特征可提升储能调度策略在已知事件场景下的收益与尾部风险表现；在先验未覆盖的负电价新事件上，松弛正则化下的语言增强策略显著优于纯数值 SAC（+5,482 元，5/5 种子），并在德国真实电价与真实气象上保持显著优势（+2,201 EUR，4/5 种子）。下一阶段将重点补强 LLM 独特性论证、完整优化基线、实测数据验证与学位论文定稿。')
    r.font.size = Pt(12); r.font.italic = True; set_cn(r, CN)

    out = os.path.join(OUT_DIR, '中期检查报告.docx')
    doc.save(out)
    return out


PPT_DARK = PptRGB(0x1F, 0x3A, 0x5F)
PPT_ACCENT = PptRGB(0x0E, 0x7C, 0x86)
PPT_WHITE = PptRGB(0xFF, 0xFF, 0xFF)
PPT_GRAY = PptRGB(0x55, 0x55, 0x55)
PPT_LIGHT = PptRGB(0xF2, 0xF5, 0xF9)
PPT_DARKT = PptRGB(0x22, 0x22, 0x22)
PPT_FONT = '微软雅黑'

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s

def slide_content(prs, title, bullets, idx, table=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    rect(slide, 0, 0, sw, Inches(0.9), PPT_DARK)
    rect(slide, 0, Inches(0.9), sw, Inches(0.06), PPT_ACCENT)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = PptPt(26); r.font.bold = True; r.font.color.rgb = PPT_WHITE; r.font.name = PPT_FONT

    top = Inches(1.2)
    if table is not None:
        rows, cols = len(table), len(table[0])
        left = Inches(0.8); width = Inches(11.4); height = Inches(0.42) * rows
        gt = slide.shapes.add_table(rows, cols, left, top, width, height).table
        for j, h in enumerate(table[0]):
            c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = PPT_DARK
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(h)
            r.font.size = PptPt(14); r.font.bold = True; r.font.color.rgb = PPT_WHITE; r.font.name = PPT_FONT
        for i in range(1, rows):
            for j in range(cols):
                c = gt.cell(i, j); c.fill.solid()
                c.fill.fore_color.rgb = PPT_LIGHT if i % 2 == 1 else PPT_WHITE
                tf = c.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
                r = p.add_run(); r.text = str(table[i][j])
                r.font.size = PptPt(13); r.font.color.rgb = PPT_DARKT; r.font.name = PPT_FONT
        top = top + height + Inches(0.15)

    if bullets:
        tb = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.0), sh - top - Inches(0.45))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = PptPt(7); p.line_spacing = 1.25
            if b.startswith('    '):
                r = p.add_run(); r.text = '· ' + b.strip()
                r.font.size = PptPt(16); r.font.color.rgb = PPT_GRAY; r.font.name = PPT_FONT
            else:
                r = p.add_run(); r.text = '• ' + b
                r.font.size = PptPt(18); r.font.color.rgb = PPT_DARKT; r.font.name = PPT_FONT

    tb = slide.shapes.add_textbox(Inches(0.3), sh - Inches(0.4), Inches(12.3), Inches(0.3))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = '中期答辩 | 李星辰 | ' + str(idx)
    r.font.size = PptPt(10); r.font.color.rgb = PPT_GRAY; r.font.name = PPT_FONT
    return slide

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = prs.slide_width, prs.slide_height
    rect(slide, 0, 0, sw, Emu(int(sh * 0.42)), PPT_DARK)
    rect(slide, 0, Emu(int(sh * 0.42)), sw, Emu(int(sh * 0.03)), PPT_ACCENT)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12.1), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = '大语言模型与深度强化学习融合的\n虚拟电厂动态优化与决策研究'
    r.font.size = PptPt(36); r.font.bold = True; r.font.color.rgb = PPT_WHITE; r.font.name = PPT_FONT
    tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.7), Inches(12.1), Inches(2.4))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for i, (k, v) in enumerate([('研究生', '李星辰'), ('导　　师', '王继军'),
                                 ('学　　校', '上海电力大学'), ('答辩日期', '2026 年 7 月 9 日')]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_after = PptPt(8)
        r = p.add_run(); r.text = k + '：' + v
        r.font.size = PptPt(20); r.font.color.rgb = PPT_GRAY; r.font.name = PPT_FONT
    return slide

def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide_cover(prs)

    slide_content(prs, '汇报提纲',
        ['一、研究背景与意义', '二、研究内容与技术路线', '三、已完成研究工作',
         '四、阶段性实验结果', '五、主要创新点', '六、存在问题与下一步计划'], 2)

    slide_content(prs, '研究背景：新型电力系统下的 VPP 调度挑战',
        ['双碳目标加新型电力系统 → 高比例新能源接入，源网荷储协同互动',
         '虚拟电厂聚合光伏、储能、可调负荷、电动汽车，参与市场与需求响应',
         'VPP 调度面临三重不确定性：负荷/光伏波动、市场价格波动、外部文本事件',
         '    气象预警、市场公告、需求响应邀约、新能源消纳提示',
         '关键痛点：文本事件先于数值效应出现，传统优化与仅数值状态 RL 难以利用'], 3)

    slide_content(prs, '问题动机：文本信息盲区与 LLM 落地风险',
        ['现有 DRL 调度多仅消费数值状态 → 高温/价格尖峰/消纳场景缺少前瞻性',
         'LLM 可理解中文事件文本，但直接让其输出调度动作存在风险：',
         '    物理可行性（SOC/功率约束）',
         '    可审计性与可复现性',
         '    市场结算逻辑一致性',
         '本研究立场：LLM 作语义编码器，不作最终决策器'], 4)

    slide_content(prs, '研究目标与科学问题',
        ['总目标：构建语言增强深度强化学习框架，提升 VPP 储能调度的事件感知能力',
         '科学问题：',
         '    如何将中文事件文本结构化为 RL 可用语义特征，且不丧失调度可复现性？',
         '    学习 Actor 与确定性语义先验如何分工？权重应否随事件覆盖变化？',
         '    当先验未覆盖的新事件类别出现时，LLM 语义特征能否帮助 SAC 适应？',
         '    上述适应性是否为仿真产物？能否在真实电价/气象下成立？'], 5)

    slide_content(prs, '技术路线',
        ['公开数据校准 → VPP 仿真环境 → 文本语义编码 → SAC 训练',
         '→ 语义安全层 → 多场景评估 → 新事件自适应 → 真实电价验证 → 学习型门控',
         '',
         '配图建议：论文总体技术架构图（fig1_overall_research_architecture）'], 6)

    slide_content(prs, '已完成工作一：公开数据校准的数据管线',
        ['国家能源局 2024 全社会用电量 98521 亿 kWh（+6.8%）',
         '可再生能源累计装机 18.89 亿 kW（风电 5.21 亿、光伏 8.87 亿）',
         '广东现货市场与独立储能公开披露（价格区间、充放电均价、价差）',
         '据此校准生成广东 VPP 15 分钟样例场景',
         '数据边界明确界定：公开数据校准背景与参数；RL 轨迹为可复现仿真序列'], 7)

    slide_content(prs, '已完成工作二：VPP 强化学习环境',
        ['状态：负荷、光伏、电价、温度、SOC、时间编码 + 文本事件语义',
         '动作：储能连续充放电功率（正放电/负充电）',
         '约束：功率边界、SOC 上下限、充放电效率',
         '奖励：购售电收益 − 退化成本 − 弃光惩罚 − 不可行惩罚',
         '统一平台：规则策略、SAC、语言增强 SAC 公平对比'], 8)

    slide_content(prs, '已完成工作三：LE-DRL-SAC 框架',
        ['DeepSeek 作语义编码器：每事件模板查询一次并缓存，输出五维语义向量',
         '    总体风险 / 价格尖峰 / 负荷压力 / 消纳压力 / 储能偏置',
         'LLM 不参与逐步控制，无在线延迟，不输出调度动作',
         '增强状态 s_aug=[s_num, s_sem] 输入 SAC',
         'Actor 语义一致性正则：训练阶段促使 Actor 学语义对应行为',
         '本地关键词编码器作 API 不可用时的可复现回退'], 9)

    slide_content(prs, '已完成工作四：语义安全层与事件覆盖权重',
        ['测试期混合：a_final=(1−w)·a_SAC + w·a_sem，再经可行性裁剪',
         '核心原则：权重 w 视事件覆盖而定，非通用常数',
         '    已知事件：先验接近场景最优 → 高 w（主实验 w=0.9）',
         '    未覆盖新事件：先验无对应分支 → 低 w，Actor 接管',
         '数值安全层消融：纯数值先验加同公式 → 隔离文本语义与机制贡献'], 10)

    slide_content(prs, '阶段性结果一：跨场景基线对比',
        ['四场景（S1 常规/S2 高温/S3 价格尖峰/S4 消纳）×3 种子 ×80 episode',
         'LE-DRL-SAC+安全层(w=0.9)：−208,468.9 元（最优）',
         '优于 SAC-Numeric(−211,023)、Rule-Based(−208,970)、增强滚动时域(−209,988)、Linear-MPC(−210,656)',
         '配对 bootstrap vs SAC-Numeric：+2,554 元，95% CI [+2,060, +3,044]，Wilcoxon p=0.0005',
         '文本消融：SAC-Numeric ≈ LE-DRL w/o Text → 增益非来自维度增大'],
        11,
        table=[['模型', '总收益(元)', 'CVaR 5%', '吞吐(MWh)', '高价放电率'],
               ['LE-DRL-SAC+安全层', '−208,468.9', '−696.7', '28.37', '0.231'],
               ['Rule-Based', '−208,969.9', '−705.0', '19.94', '0.176'],
               ['SAC-Numeric+数值安全层', '−209,125.3', '−704.0', '16.92', '0.193'],
               ['增强滚动时域', '−209,988.3', '−713.8', '14.80', '0.095'],
               ['Linear-MPC', '−210,656.0', '−714.6', '10.81', '0.073'],
               ['SAC-Numeric', '−211,023.2', '−721.3', '3.04', '0.000'],
               ['LE-DRL w/o Text', '−211,023.1', '−721.3', '3.04', '0.000']])

    slide_content(prs, '阶段性结果二：安全层权重扫描与机制隔离',
        ['w 扫描 0→1：收益单调改善，w=0.9 为保留学习分量的混合配置',
         '数值安全层：SAC-Numeric −211,023 → −209,125（机制本身有价值）',
         '语义安全层：进一步到 −208,469（文本语义额外贡献）',
         '结论：增益 = 安全层机制 + 结构化文本语义，二者可分离'], 12)

    slide_content(prs, '阶段性结果三：新事件自适应（S5 负电价）',
        ['S5：午间光伏过剩压低现货至负（最低 −210 元/MWh），晚高峰受抑',
         '先验晚高峰放电源分支失配；DeepSeek 正确识别过剩本质（消纳 0.90/偏置 +0.80）',
         '松弛正则化（λ=0, w=0）加 S1–S5 重训练，5 种子：',
         '    LE-DRL-SAC：−141,068 元 vs SAC-Numeric：−146,550 元',
         '    间隙 +5,482 元，95% CI [+4,875, +6,074]，5/5 种子有利',
         '行为：100% 负电价步充电，学会避免晚高峰放电'], 13)

    slide_content(prs, '阶段性结果四：正则双重角色与参数变体稳健性',
        ['正则扫描：已知事件收益跨权重持平（<0.5%）；S5 随松弛改善 18,515 元',
         '→ 语义信息两种可分离作用：可学状态特征 vs 约束目标，新事件上冲突',
         '参数变体 V1–V4（负电价深度/季节/光伏装机三维度）：',
         '    V1/V3/V4 显著正间隙（CI 全正，5/5 种子）',
         '    V2 趋正',
         '结论：新事件自适应非单一参数设定产物，跨维度稳健'], 14)

    slide_content(prs, '阶段性结果五：真实电价与真实气象验证',
        ['最接近真实市场的证据',
         '真实德国 DE-LU 日前价格（2026-07-04，27 个负电价区间，最低 −11.8 EUR/MWh）',
         '加 真实 Open-Meteo 气象',
         '5 种子：LE-DRL-SAC −34,956 EUR vs SAC-Numeric −37,156 EUR',
         '间隙 +2,201 EUR，95% CI [+593, +3,627]，4/5 种子有利',
         '意义：新事件自适应可迁移至真实负电价与真实气象，非仿真产物',
         '→ 为标题事件感知提供最贴近部署的实证锚点'], 15)

    slide_content(prs, '阶段性结果六：学习型事件覆盖门控',
        ['两个 SAC 专家（带/不带语义正则）加小门控网络 g_ψ(s)→w',
         '以逐批归一化 Q 差异为二分类目标训练，无需覆盖标签',
         '结果：',
         '    已知事件 S1–S4：平均 w=0.65（偏向正则化先验专家）',
         '    S5 加变体 V1–V4：w=0.31–0.48（偏向自由专家）',
         '结论：仅凭 Q 差异信号即可从数据恢复已知高 w/新事件低 w 分割 → 覆盖权重可自动化'], 16)

    slide_content(prs, '主要创新点',
        ['解耦式语言增强 SAC：LLM 严格作语义编码器（每模板查询缓存），不输出动作，决策确定可审计',
         '事件覆盖相关语义安全层：权重 w 随覆盖变化，数值安全层消融隔离语义/机制贡献',
         '新事件自适应：先验未覆盖负电价事件，松弛正则下 LLM 特征助 SAC 适应（+5,482 元，5/5）',
         '真实电价加真实气象验证：德国 DE-LU 真实价格加 Open-Meteo 真实气象（+2,201 EUR，4/5）',
         '学习型事件覆盖门控：Q 差异无标签信号恢复覆盖分割，权重可自动化'], 17)

    slide_content(prs, '存在问题与不足',
        ['数据真实性边界：主轨迹为公开数据校准仿真，非实测遥测（OOD 加真实电价已部分缓解）',
         'DeepSeek 相对关键词编码器的独特优势待专用对比坐实',
         '统计方法已修正：种子配对 bootstrap 加 Wilcoxon（旧版 3 均值 bootstrap CI 过窄已替换）',
         '完整 MILP/MPC 滚动优化基线在相同约束下对比仍需补',
         '门控仅演示覆盖分割可学，混合动作奖励最优化待做'], 18)

    slide_content(prs, '下一步研究计划',
        ['重点：补 LLM 独特性论证 / 完整优化基线 / 实测数据验证 / 门控与安全约束 / 学位论文写作'],
        19,
        table=[['阶段', '时间', '主要工作'],
               ['中期答辩', '2026.07', '完成中期检查，汇总专家意见'],
               ['实验完善', '2026.08–10', 'DeepSeek vs 关键词/noisy 消融；完整 MILP/MPC 基线；门控奖励最优化'],
               ['论文定稿', '2026.11–12', 'IEEE Access 投稿修订；学位论文各章定稿'],
               ['预答辩送审', '2027.01–03', '查重、预答辩、外审修改'],
               ['答辩', '2027 上半年', '学位论文答辩']])

    slide_content(prs, '总结与致谢',
        ['已构建语言增强 SAC 加语义安全层加新事件自适应加真实电价验证加学习型门控完整闭环',
         '核心结论：LLM 语义特征可提升 VPP 调度事件感知能力，且可迁移至真实电价',
         '核心立场：LLM 作顾问而非决策器，最终动作确定、可审计、物理可行',
         '',
         '恳请各位老师批评指正！'], 20)

    out = os.path.join(OUT_DIR, '中期答辩PPT.pptx')
    prs.save(out)
    return out


if __name__ == '__main__':
    d = build_docx(); print('DOCX:', d)
    p = build_pptx(); print('PPTX:', p)
